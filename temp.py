from pptx import Presentation
import re
import os
from tkinter import Tk
from tkinter.filedialog import askopenfilename
from docx import Document

def extract_text_and_find_latex(ppt_file):
    presentation = Presentation(ppt_file)
    id_equation_dict = {}

    with open('latex_parts.tex', 'w') as f:
        f.write("\\documentclass{article}\n\\usepackage[fleqn]{amsmath}\n\\usepackage{amssymb}\n\\setlength\\mathindent{0pt}\n\\begin{document}\n")
        id_counter = 1

        for i, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    latex_parts = re.findall(r'\\\(.*?\\\)', paragraph_text)
                    for latex in latex_parts:
                        id_text = f"**ID{id_counter:03}**"
                        paragraph_text = paragraph_text.replace(latex, id_text)
                        f.write(f"Slide {i}, {id_text}:\n\\[\n{latex[2:-2]}\n\\]\n")
                        id_counter += 1

                    for run in paragraph.runs:
                        r = run._r
                        r.getparent().remove(r)

                    new_run = paragraph.add_run()
                    new_run.text = paragraph_text

        f.write("\\end{document}\n")

    directory, filename = os.path.split(ppt_file)
    updated_filename = "updated_" + filename
    updated_file_path = os.path.join(directory, updated_filename)
    presentation.save(updated_file_path)

    pandoc_command = f'pandoc latex_parts.tex -o latex_parts.docx'
    os.system(pandoc_command)

    doc = Document('latex_parts.docx')

    current_id = None
    for paragraph in doc.paragraphs:
        parts = paragraph.text.split(':')
        if len(parts) >= 2 and "**ID" in parts[0]:
            slide_number, current_id = parts[0].split(',')
            current_id = current_id.strip()
            id_equation_dict[current_id] = parts[1].strip()
            print(f"Added to id_equation_dict: {current_id} : {id_equation_dict[current_id]}")
        elif current_id is not None:
            id_equation_dict[current_id] += ' ' + parts[0].strip()

    print(f"id_equation_dict: {id_equation_dict}")

    presentation = Presentation(updated_file_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ''.join(run.text for run in paragraph.runs)
                id_texts = re.findall(r'\*\*ID\d{3}\*\*', paragraph_text)

                for id_text in id_texts:
                    equation = id_equation_dict.get(id_text)
                    if equation is not None:
                        paragraph_text = paragraph_text.replace(id_text, equation)
                        print(f"Replaced {id_text} with {equation}")

                for run in paragraph.runs:
                    r = run._r
                    r.getparent().remove(r)

                new_run = paragraph.add_run()
                new_run.text = paragraph_text

    final_filename = "final_" + filename
    final_file_path = os.path.join(directory, final_filename)
    presentation.save(final_file_path)

if __name__ == "__main__":
    Tk().withdraw()
    ppt_file = askopenfilename(title="Select PowerPoint File",
                               filetypes=(("PowerPoint files", "*.ppt;*.pptx"), ("All files", "*.*")))
    if ppt_file:
        extract_text_and_find_latex(ppt_file)
    else:
        print("No PowerPoint file selected. Exiting...")
