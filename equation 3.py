from lxml import etree as ET

from pptx import Presentation
from pptx.oxml.ns import qn
import re
import os
import zipfile
from tkinter import Tk
from tkinter.filedialog import askopenfilename

def extract_equations_from_docx(docx_file):
    with zipfile.ZipFile(docx_file, 'r') as docx_zip:
        xml_content = docx_zip.read('word/document.xml').decode()
    xml_content = re.sub(r'<\?xml.*\?>', '', xml_content)
    root = ET.fromstring(xml_content)

    ns = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'm': 'http://schemas.openxmlformats.org/officeDocument/2006/math'
    }
    equations = []
    for oMath in root.findall('.//m:oMath', ns):
        equations.append(ET.tostring(oMath, encoding='unicode'))
    return equations

def extract_text_and_find_latex(ppt_file):
    presentation = Presentation(ppt_file)
    id_equation_dict = {}
    directory, filename = os.path.split(ppt_file)

    with open('latex_parts.tex', 'w') as f:
        f.write("\\documentclass{article}\n\\usepackage{amsmath}\n\\usepackage{amssymb}\n\\begin{document}\n")
        id_counter = 1
        for i, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    latex_parts = re.findall(r'\\\(.*?\\\)', paragraph_text)
                    for latex in latex_parts:
                        id_text = "**ID" + str(id_counter) + "**"
                        paragraph_text = paragraph_text.replace(latex, id_text)
                        f.write(f"Slide {i}, {id_text}:\n\\begin{{equation}}\n{latex[2:-2]}\n\\end{{equation}}\n")
                        id_counter += 1
                    for run in paragraph.runs:
                        r = run._r
                        r.getparent().remove(r)
                    new_run = paragraph.add_run()
                    new_run.text = paragraph_text
        f.write("\\end{document}\n")

    updated_filename = "updated_" + filename
    updated_file_path = os.path.join(directory, updated_filename)
    presentation.save(updated_file_path)

    # Debugging: Print the XML content of the updated PPTX
    with zipfile.ZipFile(updated_file_path, 'r') as ppt_zip:
        xml_content_updated = ppt_zip.read('ppt/slides/slide1.xml').decode()
        print("\n==== XML of Updated PPTX ====\n", xml_content_updated, "\n==============================\n")

    pandoc_command = f'pandoc latex_parts.tex -o latex_parts.docx'
    os.system(pandoc_command)
    equations = extract_equations_from_docx('latex_parts.docx')

    for i, equation in enumerate(equations, start=1):
        id_equation_dict["**ID" + str(i) + "**"] = equation

    print(f"Extracted equations: {id_equation_dict}")
    presentation = Presentation(updated_file_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ''.join(run.text for run in paragraph.runs)
                id_texts = re.findall(r'\*\*ID\d+\*\*', paragraph_text)
                for id_text in id_texts:
                    equation = id_equation_dict.get(id_text)
                    if equation is not None:
                        new_run = paragraph.add_run()
                        omml_r = new_run._r
                        omml_obj = ET.XML(equation)
                        omml_r.append(omml_obj)
                        paragraph_text = paragraph_text.replace(id_text, "")
                for run in paragraph.runs:
                    r = run._r
                    if not r.find(qn('m:oMath')):
                        r.getparent().remove(r)
                new_run = paragraph.add_run()
                new_run.text = paragraph_text
    final_filename = "final_" + filename
    final_file_path = os.path.join(directory, final_filename)
    presentation.save(final_file_path)

    # Debugging: Print the XML content of the final PPTX
    with zipfile.ZipFile(final_file_path, 'r') as ppt_zip:
        xml_content_final = ppt_zip.read('ppt/slides/slide1.xml').decode()
        print("\n==== XML of Final PPTX ====\n", xml_content_final, "\n===========================\n")

if __name__ == "__main__":
    Tk().withdraw()
    ppt_file = askopenfilename(title="Select PowerPoint File",
                               filetypes=(("PowerPoint files", "*.ppt;*.pptx"), ("All files", "*.*")))
    if ppt_file:
        extract_text_and_find_latex(ppt_file)
    else:
        print("No PowerPoint file selected. Exiting...")
